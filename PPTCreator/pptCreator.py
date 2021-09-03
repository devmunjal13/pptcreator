from pptx import Presentation
from pptx.util import Inches
from wand.image import Image

prs = Presentation()

def add_slide(image,titleText,subtitleText):
    lyt = prs.slide_layouts[1]
    slide = prs.slides.add_slide(lyt)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titleText
    subtitle.text = subtitleText
    imgSize = Image(filename = image)
    print(imgSize.width,imgSize.height)
    slide.shapes.add_picture(image,Inches(1),Inches(2.5),width=imgSize.width*600,height=imgSize.height*600)

def createPPT(image,titleText,subtitleText):
    add_slide(image,titleText,subtitleText)
    

def savePPT():
    prs.save("./Output/assignment.pptx")
